#!/usr/bin/env python3
# Script: build_deck.py
# Purpose: Generate the informal "Intro to Viper" slide deck.

"""Generate the informal "Intro to Viper" slide deck (16:9, dark theme).

Reproducible: edit the content/constants below and re-run:

    /tmp/pptx-venv/bin/python misc/presentation/build_deck.py

Output: misc/presentation/Viper_Intro_Revised.pptx
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --------------------------------------------------------------------------- #
# Theme tokens (GitHub-dark inspired, "viper green" accent)
# --------------------------------------------------------------------------- #
BG      = RGBColor(0x0D, 0x11, 0x17)   # slide background
PANEL   = RGBColor(0x16, 0x1B, 0x22)   # card / panel
CODEBG  = RGBColor(0x0A, 0x0E, 0x14)   # code block background
BORDER  = RGBColor(0x30, 0x36, 0x3D)   # subtle borders
TEXT    = RGBColor(0xE6, 0xED, 0xF3)   # primary text
MUTED   = RGBColor(0x8B, 0x94, 0x9E)   # secondary text
GREEN   = RGBColor(0x3F, 0xB9, 0x50)   # primary accent (viper green)
CYAN    = RGBColor(0x58, 0xA6, 0xFF)   # secondary accent (links/headers)
AMBER   = RGBColor(0xD2, 0x99, 0x22)   # highlight
COMMENT = RGBColor(0x6E, 0x96, 0x6E)   # code comments

# Syntax-highlight palette (GitHub-dark code theme). NS/PROMPT/COMMENT/TEXT
# reuse the theme tokens above so the code blocks stay on-palette.
SYN_KW  = RGBColor(0xFF, 0x7B, 0x72)   # keywords / operators
SYN_FN  = RGBColor(0xD2, 0xA8, 0xFF)   # function / method / command names
SYN_STR = RGBColor(0xA5, 0xD6, 0xFF)   # strings / paths / urls
SYN_NUM = RGBColor(0x79, 0xC0, 0xFF)   # numbers / flags
SYN_NS  = GREEN                        # namespace / type segments
SYN_PR  = MUTED                        # shell prompts
SYN_CMT = COMMENT                      # comments

SANS = "Avenir Next"
MONO = "Menlo"

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, "..", ".."))
LOGO = os.path.join(ROOT, "misc", "images", "viperlogo2.png")
OUT  = os.path.join(HERE, "Viper_Intro_Revised.pptx")

EMU_IN = 914400
SW, SH = 13.333, 7.5   # widescreen inches

prs = Presentation()
prs.slide_width = Emu(int(SW * EMU_IN))
prs.slide_height = Emu(int(SH * EMU_IN))
BLANK = prs.slide_layouts[6]

# Slides that carry no page-number footer (title, section dividers, closing).
# Tracked by 1-based index because python-pptx may hand back fresh Slide
# wrappers on each access, so object identity is not stable.
NO_FOOTER = []


# --------------------------------------------------------------------------- #
# Low-level helpers
# --------------------------------------------------------------------------- #
def _no_shadow(shape):
    shape.shadow.inherit = False


def new_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    _no_shadow(bg)
    return slide


def no_footer(slide):
    """Mark the just-created slide as having no page-number footer."""
    NO_FOOTER.append(len(prs.slides))
    return slide


def rect(slide, l, t, w, h, fill=PANEL, line=BORDER, lw=1.0,
         rounded=True, radius=0.11):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    _no_shadow(sp)
    if rounded:
        try:
            # radius given in inches → fraction of the shorter side
            sp.adjustments[0] = max(0.0, min(0.5, radius / min(w, h)))
        except Exception:
            pass
    return sp


def down_arrow(slide, cx, t, color=BORDER, w=0.28, h=0.26):
    sp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                Inches(cx - w / 2), Inches(t),
                                Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    _no_shadow(sp)
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def para(tf, runs, size=18, color=TEXT, bold=False, font=SANS,
         align=PP_ALIGN.LEFT, before=0, after=6, line=1.05, first=False,
         bullet=None):
    """Add a paragraph. `runs` is a str or list of (text, overrides) tuples."""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else \
        tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    try:
        p.line_spacing = line
    except Exception:
        pass
    if isinstance(runs, str):
        runs = [(runs, {})]
    if bullet is not None:
        runs = [(bullet + "  ", {"color": GREEN, "bold": True})] + list(runs)
    for text, ov in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(ov.get("size", size))
        r.font.bold = ov.get("bold", bold)
        r.font.name = ov.get("font", font)
        r.font.color.rgb = ov.get("color", color)
    return p


# --------------------------------------------------------------------------- #
# Syntax highlighting
#
# Two small tokenizers reproduce the deck's GitHub-dark code styling:
#   highlight_code    — Zia / IL / BASIC (keyword + dotted-namespace rules)
#   highlight_console — shell / REPL transcripts (command + path + prompt)
# Each returns a list of (text, RGBColor) spans; run boundaries are not
# significant, only the per-character colour.
# --------------------------------------------------------------------------- #
CODE_KW = {
    # Zia
    "module", "bind", "as", "func", "var",
    # IL
    "il", "global", "const", "str", "void", "alloca", "load", "store",
    "call", "ret",
    # BASIC
    "USING", "LET", "PRINT",
}
CODE_OPS = ("->", "=", "+", "*", "%")          # match "->" before "-"
QNAME = re.compile(r"[A-Z]\w*(?:\.[A-Za-z_]\w*)+")  # Title-cased dotted chain

CONSOLE_CMDS = ("viper", "git", "cd", "clone", "init", "run", "build",
                "repl", "il-opt", "make", "Say")


def _span(spans, text, color):
    if text:
        spans.append((text, color))


def highlight_code(line):
    """Tokenise a Zia/IL/BASIC line into coloured spans."""
    spans = []
    n = len(line)
    # A leading "label:" (IL block label) colours the label as a keyword.
    lbl = re.match(r"(\s*)([A-Za-z_]\w*):", line)
    label_end = lbl.end(2) if lbl else -1
    i = 0
    pending_ns = False        # identifier right after `as` is a namespace alias
    while i < n:
        c = line[i]
        # capture-and-clear so only whitespace preserves the alias flag
        alias, pending_ns = pending_ns, False
        if c in " \t":
            j = i
            while j < n and line[j] in " \t":
                j += 1
            _span(spans, line[i:j], TEXT)
            pending_ns = alias
            i = j
            continue
        # comments: REM (BASIC), '#', '//'
        if line.startswith("REM", i) and (i + 3 == n or not line[i + 3].isalnum()):
            _span(spans, line[i:], SYN_CMT)
            break
        if c == "#" or line.startswith("//", i):
            _span(spans, line[i:], SYN_CMT)
            break
        if c == '"':
            j = line.find('"', i + 1)
            j = j + 1 if j != -1 else n
            _span(spans, line[i:j], SYN_STR)
            i = j
            continue
        if i < label_end:                       # the IL label word
            _span(spans, line[lbl.start(2):label_end], SYN_KW)
            i = label_end
            continue
        m = QNAME.match(line, i)
        if m:
            qn = m.group(0)
            after = i + len(qn)
            is_call = after < n and line[after] == "("
            segs = qn.split(".")
            for k, seg in enumerate(segs):
                if k:
                    _span(spans, ".", TEXT)
                leaf = (k == len(segs) - 1)
                _span(spans, seg, SYN_FN if (leaf and is_call) else SYN_NS)
            i = after
            continue
        m = re.match(r"[A-Za-z_]\w*", line[i:])
        if m:
            w = m.group(0)
            if alias:
                col = SYN_NS
            elif w in CODE_KW:
                col = SYN_KW
            else:
                col = TEXT
            _span(spans, w, col)
            if w == "as":
                pending_ns = True
            i += len(w)
            continue
        m = re.match(r"\d+", line[i:])
        if m:
            _span(spans, m.group(0), SYN_NUM)
            i += len(m.group(0))
            continue
        for op in CODE_OPS:
            if line.startswith(op, i):
                _span(spans, op, SYN_KW)
                i += len(op)
                break
        else:
            _span(spans, c, TEXT)
            i += 1
    return spans


_CONSOLE_PATH = (
    r"(?<!\S)(?:\./|\.\\|https?://)\S+"            # ./path, urls
    r"|(?<!\S)[\w.~-]+(?:/[\w.~/-]+)+"             # a/b/c paths
    r"|(?<!\S)[A-Za-z0-9_~-]+\.[A-Za-z]\w*"        # file.ext
)


def highlight_console(line, numbers=True):
    """Tokenise a shell/REPL line into coloured spans."""
    spans = []
    start = 0
    pm = re.match(r"\s*(?:\$|zia>|>>>)", line)
    if pm:
        _span(spans, line[:pm.end()], SYN_PR)
        start = pm.end()
    rest = line[start:]
    alts = [
        r"(?P<comment>\#.*$)",
        r'(?P<string>"[^"]*")',
        r"(?P<path>" + _CONSOLE_PATH + r")",
        r"(?P<flag>(?<!\S)-\w+)",
        r"(?P<cmd>\b(?:" + "|".join(CONSOLE_CMDS) + r")\b)",
    ]
    if numbers:
        alts.append(r"(?P<number>(?<!\w)\d+(?=\s|$))")
    rx = re.compile("|".join(alts), re.M)
    colors = {"comment": SYN_CMT, "string": SYN_STR, "path": SYN_STR,
              "flag": SYN_NUM, "cmd": SYN_FN, "number": SYN_NUM}
    pos = 0
    for m in rx.finditer(rest):
        if m.start() > pos:
            _span(spans, rest[pos:m.start()], TEXT)
        _span(spans, m.group(), colors[m.lastgroup])
        pos = m.end()
    if pos < len(rest):
        _span(spans, rest[pos:], TEXT)
    return spans


def hl_para(tf, spans, size, first=False, after=2, line=1.0):
    """Add a code paragraph from highlighter spans (monospace runs)."""
    runs = [(t, {"color": c, "font": MONO}) for t, c in spans]
    if not runs:
        runs = [(" ", {"color": TEXT, "font": MONO})]
    return para(tf, runs, size=size, first=first, after=after, line=line)


def kicker(slide, text):
    """Small green eyebrow label + a short accent rule above the title."""
    tf = textbox(slide, 0.7, 0.42, 11.9, 0.4)
    para(tf, text.upper(), size=13, color=GREEN, bold=True, first=True,
         after=0)
    rect(slide, 0.72, 0.86, 0.55, 0.055, fill=GREEN, line=None, rounded=False)


def title(slide, text, sub=None, th=1.0):
    tf = textbox(slide, 0.7, 0.98, 11.9, th)
    para(tf, text, size=34, color=TEXT, bold=True, first=True, after=2)
    if sub:
        para(tf, sub, size=16, color=MUTED, after=0)


def code_block(slide, l, t, w, h, lines, header=None, size=13.5,
               lang="code", header_color=MUTED):
    """A syntax-highlighted code panel.

    lang: "code" (Zia/IL/BASIC), "console" (shell with numbers) or
    "box" (command list, no numeric highlighting).
    """
    if header:
        htf = textbox(slide, l + 0.02, t - 0.34, w, 0.3)
        para(htf, header, size=12, color=header_color, bold=True, first=True,
             after=0)
    panel = rect(slide, l, t, w, h, fill=CODEBG, line=BORDER, lw=1.0,
                 radius=0.1)
    tf = panel.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.14)
    tf.margin_bottom = Inches(0.10)
    for i, raw in enumerate(lines):
        if lang == "console":
            spans = highlight_console(raw, numbers=True)
        elif lang == "box":
            spans = highlight_console(raw, numbers=False)
        else:
            spans = highlight_code(raw)
        hl_para(tf, spans, size=size, first=(i == 0), after=2, line=1.0)
    return panel


def card(slide, l, t, w, h, accent=BORDER, lw=1.5):
    """Rounded panel with an accent-colored border (no stripe); returns it."""
    return rect(slide, l, t, w, h, fill=PANEL, line=accent, lw=lw,
                radius=0.13)


def card_text(card_shape, ml=0.3, mt=0.27, mr=0.28, mb=0.2):
    """Configure a card's own text frame for top-anchored, wrapped body text."""
    tf = card_shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)
    return tf


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #
def slide_title():
    s = new_slide()
    no_footer(s)
    # left accent column
    rect(s, 0, 0, 0.16, SH, fill=GREEN, line=None, rounded=False)
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, Inches(0.95), Inches(1.05),
                             height=Inches(1.9))
    tf = textbox(s, 0.95, 3.2, 11.5, 2.6)
    para(tf, "An informal introduction", size=16, color=GREEN, bold=True,
         first=True, after=6)
    para(tf, "Viper", size=76, color=TEXT, bold=True, after=2)
    para(tf, "A from-scratch compiler, VM, runtime, and native toolchain "
            "for games and apps.", size=22, color=MUTED,
         after=0, line=1.1)
    # footer chips
    foot = textbox(s, 0.95, 6.7, 11.5, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    para(foot, [
        ("Pre-alpha / active development", {"color": TEXT, "bold": True}),
        ("    ·    ", {"color": BORDER}),
        ("GPL-3.0", {"color": MUTED}),
        ("    ·    ", {"color": BORDER}),
        ("Linux · macOS · Windows", {"color": MUTED}),
        ("    ·    ", {"color": BORDER}),
        ("github.com/splanck/viper", {"color": CYAN}),
    ], size=14, first=True, after=0)


def slide_origin():
    s = new_slide()
    no_footer(s)
    rect(s, 0, 0, 0.16, SH, fill=GREEN, line=None, rounded=False)
    tf = textbox(s, 0.95, 2.13, 11.5, 0.66)
    para(tf, "THE ORIGIN STORY", size=15, color=GREEN, bold=True, first=True,
         after=0)
    t2 = textbox(s, 0.95, 2.62, 11.5, 2.08)
    para(t2, "Where did Viper come from?", size=54, color=TEXT, bold=True,
         first=True, after=0)
    t3 = textbox(s, 0.95, 4.87, 10.39, 0.98)
    para(t3, "How a weekend compiler experiment quietly turned into an "
             "entire stack.", size=20, color=MUTED, first=True, after=0,
         line=1.15)


def slide_what_is():
    s = new_slide()
    kicker(s, "What is Viper")
    title(s, "A platform for building games", th=0.77)

    tf = textbox(s, 0.7, 1.91, 11.26, 1.64)
    para(tf, [
        ("Viper is a ", {"color": TEXT}),
        ("game development platform and runtime",
         {"color": GREEN, "bold": True}),
        (". It is general purpose, but the whole stack is shaped first and "
         "foremost around ", {"color": TEXT}),
        ("PC game development", {"color": CYAN, "bold": True}),
        (".", {"color": TEXT}),
    ], size=24, first=True, after=0, line=1.12)

    cards = [
        (0.7, "Games first",
         "Graphics, audio, input, 2D and 3D game systems are the core of "
         "the runtime."),
        (4.78, "General purpose",
         "The same compiler, languages, and runtime handle ordinary apps "
         "and tools too."),
        (8.86, "One whole stack",
         "Compiler, VM, native backends, and runtime are all part of Viper "
         "itself."),
    ]
    for x, head, body in cards:
        rect(s, x, 4.21, 3.78, 1.91, fill=CODEBG, line=BORDER, lw=1.0,
             radius=0.11)
        ctf = textbox(s, x + 0.25, 4.56, 3.28, 1.37)
        para(ctf, head, size=17, color=GREEN, bold=True, first=True, after=6)
        para(ctf, body, size=14, color=MUTED, after=0, line=1.18)


def slide_why():
    s = new_slide()
    kicker(s, "Why Viper")
    title(s, "The IL is the protagonist",
          "One typed, SSA-based intermediate language sits at the center "
          "of everything.")

    tf = textbox(s, 0.7, 2.3, 11.9, 1.6)
    para(tf, [
        ("Most toolchains bolt a language onto someone else's stack. Viper "
         "inverts that: a small, ", {"color": TEXT}),
        ("typed, verifiable IL", {"color": GREEN, "bold": True}),
        (" is the contract. Every language compiles to it, every backend "
         "consumes it, and one runtime serves them all.", {"color": TEXT}),
    ], size=20, first=True, after=0, line=1.32)

    # thin-waist glyph (the deck's single home for this visual)
    gy, gh = 4.4, 1.2
    bw, gap = 3.3, 0.62
    x0 = (SW - (bw * 3 + gap * 2)) / 2          # centered
    x1 = x0 + bw + gap
    x2 = x1 + bw + gap

    def node(x, label, sub, filled=False):
        if filled:
            c = rect(s, x, gy, bw, gh, fill=GREEN, line=None, radius=0.12)
            lc, sc, lsz = CODEBG, CODEBG, 20
        else:
            c = card(s, x, gy, bw, gh, accent=BORDER)
            lc, sc, lsz = TEXT, MUTED, 18
        ctf = c.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.margin_left = Inches(0.2); ctf.margin_right = Inches(0.2)
        para(ctf, label, size=lsz, color=lc, bold=True,
             align=PP_ALIGN.CENTER, first=True, after=2)
        para(ctf, sub, size=12.5, color=sc, align=PP_ALIGN.CENTER, after=0)

    def glyph_arrow(x):
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x),
                               Inches(gy + gh / 2 - 0.18), Inches(gap),
                               Inches(0.36))
        a.fill.solid(); a.fill.fore_color.rgb = GREEN
        a.line.fill.background(); _no_shadow(a)

    node(x0, "Zia · BASIC", "languages")
    glyph_arrow(x0 + bw)
    node(x1, "Viper IL", "typed SSA · verified", filled=True)
    glyph_arrow(x1 + bw)
    node(x2, "VM · Native", "run or ship")

    cap = textbox(s, 0.7, gy + gh + 0.34, 11.9, 0.5,
                  anchor=MSO_ANCHOR.MIDDLE)
    para(cap, "Frontends and backends are interchangeable — the IL is the "
              "one thing they all agree on.", size=15, color=MUTED,
         align=PP_ALIGN.CENTER, first=True, after=0)


def slide_core_idea():
    s = new_slide()
    kicker(s, "The core idea")
    title(s, "One IL. Many languages. VM and native targets.")

    tf = textbox(s, 0.7, 2.45, 6.85, 4.4)
    pts = [
        ("Programs in ", ("Zia", GREEN), " or ", ("BASIC", GREEN),
         " compile to ", ("Viper IL", GREEN),
         " — a typed, SSA-based intermediate language."),
        ("The IL is the ", ("“thin waist”", TEXT),
         ": frontends and backends plug into it independently."),
        ("Run it on the ", ("VM", CYAN),
         " for fast iteration, or compile straight to ",
         ("native machine code", CYAN), "."),
    ]
    for i, parts in enumerate(pts):
        runs = []
        for part in parts:
            if isinstance(part, tuple):
                runs.append((part[0], {"color": part[1], "bold": True}))
            else:
                runs.append((part, {"color": TEXT}))
        para(tf, runs, size=21, first=(i == 0), after=22, line=1.14,
             bullet="▸")

    # "self-contained stack" emphasis card on the right
    ctf = card_text(card(s, 7.85, 2.25, 4.78, 4.13, accent=GREEN),
                    ml=0.32, mt=0.3)
    para(ctf, "Self-contained stack", size=21, color=TEXT, bold=True,
         first=True, after=8)
    para(ctf, "Every layer of the stack is implemented in the Viper repo.",
         size=15, color=MUTED, after=16, line=1.18)
    for feat in ["Compiler, VM, assembler, linker, runtime in-tree",
                 "Native code generation without LLVM",
                 "Custom runtime and graphics layers",
                 "Minimal third-party dependencies"]:
        para(ctf, feat, size=15.5, color=TEXT, after=10, bullet="✓")


def slide_languages():
    s = new_slide()
    kicker(s, "Frontends")
    title(s, "Two languages, one target")

    cards = [
        (0.7, GREEN, "Zia", "the main language",
         "Modern, statically typed — designed for apps and games.",
         ["Classes, generics & enums", "Lambdas & modules",
          "Pattern matching & exhaustiveness",
          "Memory-safe user surface — no raw pointers in normal code"]),
        (6.78, CYAN, "BASIC", "the teaching / prototyping frontend",
         "A classic, approachable dialect for rapid prototyping.",
         ["Familiar LET / PRINT / IF / WHILE",
          "Lexer + recursive-descent parser",
          "Intrinsics lower to runtime calls",
          "Great for learning the pipeline"]),
    ]
    for left, accent, name, tag, desc, feats in cards:
        tf = card_text(card(s, left, 2.3, 5.85, 3.7, accent=accent))
        para(tf, [(name, {"color": TEXT, "bold": True, "size": 26}),
                  ("   " + tag, {"color": accent, "size": 14})],
             first=True, after=10)
        para(tf, desc, size=16, color=MUTED, after=14, line=1.12)
        for feat in feats:
            para(tf, feat, size=16, color=TEXT, after=9, bullet="✓")

    foot = textbox(s, 0.7, 6.4, 11.9, 0.45)
    para(foot, [("Both lower to the same Viper IL",
                 {"color": TEXT, "bold": True}),
                (" — so they share the same optimizer, backends, and "
                 "runtime library.", {"color": MUTED})],
         size=15, first=True, after=0)


def slide_same_program():
    s = new_slide()
    kicker(s, "Frontends")
    title(s, "Same program, two languages", th=0.77)

    zia = [
        "module Area;",
        "",
        "bind Viper.Terminal as Terminal;",
        "bind Viper.Text.Fmt as Fmt;",
        "",
        "func start() {",
        "    var w = 4;",
        "    var h = 3;",
        "    var area = w * h;",
        "    Terminal.Say(\"Area = \");",
        "    Terminal.Say(Fmt.Int(area));",
        "}",
    ]
    basic = [
        "REM Area",
        "",
        "USING Viper.Text",
        "",
        "LET W = 4",
        "LET H = 3",
        "LET AREA = W * H",
        "PRINT \"Area = \"",
        "PRINT Fmt.Int(AREA)",
    ]
    code_block(s, 0.7, 2.05, 5.55, 3.65, zia, header="Zia", size=13,
               lang="code", header_color=GREEN)
    code_block(s, 7.1, 2.05, 5.55, 3.65, basic, header="Viper BASIC",
               size=13, lang="code", header_color=CYAN)

    eq = textbox(s, 6.35, 3.48, 0.66, 0.66, anchor=MSO_ANCHOR.MIDDLE)
    para(eq, "=", size=40, color=GREEN, bold=True, align=PP_ALIGN.CENTER,
         first=True, after=0, line=1.0)

    foot = textbox(s, 0.7, 5.86, 11.9, 0.7)
    para(foot, [("Same runtime, either way", {"color": GREEN, "bold": True}),
                (" — both frontends call the very same Viper.Text.Fmt and "
                 "terminal classes. Only the syntax differs.",
                 {"color": MUTED})],
         size=15, first=True, after=0, line=1.1)


def slide_il():
    s = new_slide()
    kicker(s, "Viper IL")
    title(s, "Typed, SSA-based, and inspectable")

    zia = [
        "module Hello;",
        "",
        "bind Viper.Terminal as Terminal;",
        "bind Viper.Text.Fmt as Fmt;",
        "",
        "func start() {",
        "    var x = 2 + 3;",
        "    var y = x * 2;",
        "    Terminal.Say(\"HELLO\");",
        "    Terminal.Say(Fmt.Int(y));",
        "}",
    ]
    il = [
        "il 0.3.0",
        "global const str @.L0 = \"HELLO\"",
        "func @main() -> void {",
        "entry_0:",
        "  %t0 = iadd.ovf 2, 3",
        "  %t1 = alloca 8",
        "  store i64, %t1, %t0",
        "  %t2 = load i64, %t1",
        "  %t3 = imul.ovf %t2, 2",
        "  %t4 = alloca 8",
        "  store i64, %t4, %t3",
        "  %t5 = const_str @.L0",
        "  call @Viper.Terminal.Say(%t5)",
        "  %t6 = load i64, %t4",
        "  %t7 = call @Viper.Text.Fmt.Int(%t6)",
        "  call @Viper.Terminal.Say(%t7)",
        "  ret",
        "}",
    ]
    code_block(s, 0.7, 2.05, 5.55, 4.2, zia, header="Zia source", size=13.5,
               lang="code")
    # arrow between
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.42), Inches(3.9),
                            Inches(0.5), Inches(0.42))
    ar.fill.solid(); ar.fill.fore_color.rgb = GREEN
    ar.line.fill.background(); _no_shadow(ar)
    code_block(s, 7.1, 2.05, 5.55, 4.2, il, header="Viper IL — unoptimized",
               size=11.5, lang="code")

    foot = textbox(s, 0.7, 6.4, 11.9, 0.7)
    para(foot, [
        ("A verifier", {"color": CYAN, "bold": True}),
        (" enforces structural and type invariants, then a ",
         {"color": MUTED}),
        ("24-pass optimizer", {"color": CYAN, "bold": True}),
        (" runs GVN, LICM, SCCP, inlining, loop opts, and more.",
         {"color": MUTED}),
    ], size=15, first=True, after=0, line=1.1)


def slide_architecture():
    s = new_slide()
    kicker(s, "Architecture")
    title(s, "The end-to-end pipeline")

    cx = 6.667

    def box(t, w, h, lines, accent=BORDER, fill=PANEL, big=False, left=None):
        l = (cx - w / 2) if left is None else left
        cd = rect(s, l, t, w, h, fill=fill, line=accent,
                  lw=1.6 if accent != BORDER else 1.0, radius=0.11)
        tf = cd.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        for i, (txt, sz, col, bold) in enumerate(lines):
            para(tf, txt, size=sz, color=col, bold=bold,
                 align=PP_ALIGN.CENTER, first=(i == 0), after=0, line=1.0)
        return cd

    # uniform two-line boxes: bold heading + muted detail (no inline dashes)
    box(1.69, 5.6, 0.82,
        [("Source languages", 17, TEXT, True),
         ("Zia · BASIC", 12.5, MUTED, False)], accent=CYAN)
    down_arrow(s, cx, 2.59)
    box(2.83, 5.6, 0.82,
        [("Frontend", 17, TEXT, True),
         ("lex · parse · sema · lower", 12.5, MUTED, False)])
    down_arrow(s, cx, 3.73)
    box(3.97, 6.6, 0.88,
        [("Viper IL", 19, GREEN, True),
         ("typed SSA · verifier · 24-pass optimizer", 12.5, MUTED, False)],
        accent=GREEN)

    lx, rx = 3.97, 9.36   # centers of the split boxes
    down_arrow(s, lx, 4.93)
    down_arrow(s, rx, 4.93)
    box(5.17, 4.55, 0.94,
        [("Bytecode VM", 17, TEXT, True),
         ("switch · table · threaded dispatch", 12, MUTED, False)],
        accent=AMBER, left=1.7)
    box(5.17, 4.55, 0.94,
        [("Native path", 17, TEXT, True),
         ("AArch64 · x86-64 · assembler · linker", 12, MUTED, False)],
        accent=AMBER, left=7.08)
    down_arrow(s, lx, 6.15)
    down_arrow(s, rx, 6.15)

    box(6.41, 6.6, 0.82,
        [("Viper Runtime", 17, TEXT, True),
         ("graphics · GUI · 3D · audio · net · crypto · game systems · …",
          12.5, MUTED, False)],
        accent=CYAN)


def slide_execution_vm():
    s = new_slide()
    kicker(s, "Execution · 1 of 2")
    title(s, "Fast iteration on the VM",
          "The primary development and debugging target.")

    tf = card_text(card(s, 0.7, 2.5, 5.95, 3.75, accent=CYAN))
    para(tf, "The VM", size=24, color=CYAN, bold=True, first=True, after=12)
    for feat in ["Register-file bytecode interpreter",
                 "Switch · table · threaded-goto dispatch",
                 "Pooled frames · cached string literals",
                 "Structured trap diagnostics · execution tracing"]:
        para(tf, feat, size=16, color=TEXT, after=12, bullet="▸")

    session = [
        "$ viper run game.zia      # no build step",
        "$ viper repl              # live evaluation",
        "zia> Say(\"hello\")",
        "hello",
        "zia> Say(Fmt.Int(2 + 3))",
        "5",
    ]
    code_block(s, 7.0, 2.85, 5.6, 2.55, session,
               header="Edit → run → repeat", size=13, lang="console")
    cap = textbox(s, 7.0, 5.62, 5.6, 0.7)
    para(cap, [("No compile step", {"color": TEXT, "bold": True}),
               (" — see results right away, then build the same IL "
                "as a native binary.", {"color": MUTED})],
         size=14.5, first=True, after=0, line=1.15)


def slide_execution_native():
    s = new_slide()
    kicker(s, "Execution · 2 of 2")
    title(s, "Compile straight to a native binary",
          "Same IL, no VM required.")

    ltf = card_text(card(s, 0.7, 2.5, 5.95, 3.75, accent=GREEN))
    para(ltf, "Native path", size=24, color=GREEN, bold=True, first=True,
         after=12)
    for feat, emph in [
        ("AArch64 and x86-64 targets for macOS, Windows, and Linux", False),
        ("Register coalescing · post-RA scheduling", False),
        ("Built-in assembler + linker", True),
        ("ELF · Mach-O · PE output with debug info", False),
    ]:
        para(ltf, feat, size=16, color=TEXT, bold=emph, after=12, bullet="▸")

    rtf = card_text(card(s, 6.85, 2.5, 5.78, 3.75, accent=CYAN))
    para(rtf, "Verified equivalence", size=24, color=CYAN, bold=True,
         first=True, after=10)
    para(rtf, "VM and native outputs must match for every defined "
              "program.", size=15.5, color=MUTED, after=16, line=1.18)
    for feat in ["Differential testing — VM vs. native",
                 "~18 fuzz harnesses on parser & IL",
                 "ASAN · UBSAN · TSAN sanitizer lanes"]:
        para(rtf, feat, size=16, color=TEXT, after=12, bullet="▸")


def slide_runtime():
    s = new_slide()
    kicker(s, "The runtime")
    title(s, "One standard library, every frontend",
          "436 classes across 22 modules — shared through a C ABI by "
          "both the VM and native code.")

    modules = [
        ("Graphics", 57), ("GUI", 53), ("Graphics3D", 46), ("Game3D", 34),
        ("Collections", 29), ("Network", 27), ("Game", 26), ("Utilities", 23),
        ("Text", 22), ("Threads", 18), ("I/O", 16), ("Math", 12),
        ("Localization", 10), ("Crypto", 8), ("Sound", 8), ("Time", 8),
    ]
    cols, rows = 4, 4
    L0, T0 = 0.7, 2.55
    gx, gy = 0.22, 0.2
    cw = (12.63 - 0.7 - (cols - 1) * gx) / cols
    ch = 0.84
    for i, (name, count) in enumerate(modules):
        r, c = divmod(i, cols)
        l = L0 + c * (cw + gx)
        t = T0 + r * (ch + gy)
        cd = rect(s, l, t, cw, ch, fill=PANEL, line=BORDER, lw=1.0,
                  radius=0.11)
        tf = cd.text_frame
        tf.margin_left = Inches(0.16); tf.margin_top = Inches(0.06)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, [(str(count), {"color": GREEN, "bold": True, "size": 22}),
                  ("  classes", {"color": MUTED, "size": 10.5})],
             first=True, after=1)
        para(tf, name, size=13.5, color=TEXT, bold=True, after=0)

    foot = textbox(s, 0.7, 6.74, 11.9, 0.5)
    para(foot, [
        ("Plus", {"color": MUTED}),
        (" Core, Data, Input, Memory, Compress, Game.Physics2D, Game.UI",
         {"color": TEXT}),
        ("  —  graphics, 3D, audio, networking, crypto, "
         "threading, localization, and more.", {"color": MUTED}),
    ], size=14, first=True, after=0, line=1.1)


def slide_by_numbers():
    s = new_slide()
    kicker(s, "By the numbers")
    title(s, "What “from scratch” really means",
          "Every layer is in-tree. A sense of the scale:")

    stats = [
        ("~700K", "lines of code"),
        ("1,700+", "tests"),
        ("436", "runtime classes"),
        ("23", "apps & games"),
    ]
    T0 = 3.28
    cw, ch = 2.82, 1.6
    xs = [0.7, 3.74, 6.78, 9.81]
    for i, (num, label) in enumerate(stats):
        l = xs[i]
        cd = rect(s, l, T0, cw, ch, fill=PANEL, line=BORDER, lw=1.0,
                  radius=0.11)
        tf = cd.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, num, size=36, color=GREEN, bold=True, align=PP_ALIGN.CENTER,
             first=True, after=3, line=1.0)
        para(tf, label, size=13, color=MUTED, align=PP_ALIGN.CENTER, after=0,
             line=1.08)

    foot = textbox(s, 0.7, 6.72, 11.9, 0.5)
    para(foot, [("Pre-Alpha, under active development",
                 {"color": TEXT, "bold": True}),
                (" — measures of scope, not maturity.", {"color": MUTED})],
         size=14.5, first=True, after=0)


def slide_enough():
    s = new_slide()
    no_footer(s)
    rect(s, 0, 0, 0.16, SH, fill=GREEN, line=None, rounded=False)
    tf = textbox(s, 0.95, 2.52, 11.5, 0.66)
    para(tf, "ENOUGH ARCHITECTURE", size=15, color=GREEN, bold=True,
         first=True, after=0)
    t2 = textbox(s, 0.95, 3.01, 11.5, 1.64)
    para(t2, "Let's look at the demos.", size=60, color=TEXT, bold=True,
         first=True, after=0)
    t3 = textbox(s, 0.95, 4.7, 10.39, 0.98)
    para(t3, "Games and apps built entirely in Zia, running on the "
             "toolchain we just walked through.", size=20, color=MUTED,
         first=True, after=0, line=1.15)


def slide_showcase():
    s = new_slide()
    kicker(s, "Built with Viper")
    title(s, "Larger demos, all in pure Zia",
          "All open source in the repo — running on the same toolchain.")

    demos = [
        ("XENOSCAPE", "~21K LOC",
         "Metroidvania-style demo — bosses, abilities, saves"),
        ("ViperSQL", "~85K LOC",
         "Experimental SQL database engine + client"),
        ("Chess", "Zia", "Alpha-beta AI, drag-and-drop GUI"),
        ("3D Bowling", "Zia", "Physics-driven, multi-mode camera"),
        ("Ridgebound", "Zia", "Terrain · water · PBR · post-FX"),
        ("Crackman", "Zia", "Maze chase with ghost AI"),
        ("Paint", "Zia", "Layers, undo/redo, full tool set"),
        ("+ more", "23 total", "6 apps · 17 games/demos"),
    ]
    cols = 4
    L0, T0 = 0.7, 2.7
    gx, gy = 0.22, 0.24
    cw = (12.63 - 0.7 - (cols - 1) * gx) / cols
    ch = 1.62
    for i, (name, meta, desc) in enumerate(demos):
        r, c = divmod(i, cols)
        l = L0 + c * (cw + gx)
        t = T0 + r * (ch + gy)
        marquee = (i == 0)
        cd = card(s, l, t, cw, ch, accent=GREEN if marquee else BORDER,
                  lw=1.5 if marquee else 1.0)
        tf = cd.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.16)
        tf.margin_right = Inches(0.16)
        para(tf, name, size=16, color=TEXT, bold=True, first=True, after=3)
        para(tf, meta, size=12, color=GREEN, bold=True, after=5)
        para(tf, desc, size=12, color=MUTED, after=0, line=1.12)

    foot = textbox(s, 0.7, 6.72, 11.9, 0.45)
    para(foot, [("Same compiler, VM, and runtime", {"color": TEXT,
                 "bold": True}),
                (" — demos build through the native toolchain for the "
                 "supported desktop platforms.", {"color": MUTED})],
         size=14.5, first=True, after=0)


def slide_tooling():
    s = new_slide()
    kicker(s, "Tooling & getting started")
    title(s, "From source to native binary")

    cmds = [
        "git clone https://github.com/splanck/viper",
        "cd viper",
        "./scripts/build_viper_linux.sh    # build the toolchain",
        "./scripts/build_demos_linux.sh    # build the demos",
        "",
        "viper init my-app",
        "viper run my-app",
        "viper repl",
        "viper build my-app -o my-app",
    ]
    code_block(s, 0.7, 2.25, 6.55, 2.45, cmds, header="Build from source",
               size=12, lang="console")
    req = textbox(s, 0.72, 4.78, 6.5, 0.32)
    para(req, [("Requirements:  ", {"color": MUTED, "bold": True}),
               ("supported OS · C++ compiler · CMake · make",
                {"color": MUTED})], size=12.5, first=True, after=0)

    # tools card (right column, full height) — console-styled command list
    box = card(s, 7.5, 2.25, 5.13, 4.15, accent=GREEN)
    ttf = card_text(box, ml=0.3, mt=0.24)
    ttf.word_wrap = False
    para(ttf, "In the box", size=18, color=TEXT, bold=True, first=True,
         after=10)
    for cmd in [
        "viper  unified driver — run/build/package",
        "repl  live Zia & BASIC evaluation",
        "viper il-opt  optimize IL (24 passes)",
        "il-verify · il-dis  verify & disassemble IL",
        "zia-server  LSP + MCP language server",
        "package  installers: .app · .deb · .exe · .tar.gz",
    ]:
        hl_para(ttf, highlight_console(cmd, numbers=False), size=13.5,
                after=10, line=1.05)

    foot = textbox(s, 0.7, 6.6, 11.9, 0.45, anchor=MSO_ANCHOR.MIDDLE)
    para(foot, [("github.com/splanck/viper", {"color": CYAN, "bold": True}),
                ("     —     clone it, run the demos, and tell me what "
                 "breaks.", {"color": MUTED})],
         size=15, first=True, after=0)


def slide_thanks():
    s = new_slide()
    no_footer(s)
    rect(s, 0, 0, 0.16, SH, fill=GREEN, line=None, rounded=False)
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, Inches(0.95), Inches(1.05),
                             height=Inches(1.9))
    t1 = textbox(s, 0.95, 3.55, 11.5, 1.31)
    para(t1, "Thanks for watching.", size=60, color=TEXT, bold=True,
         first=True, after=0)
    t2 = textbox(s, 0.95, 4.76, 10.72, 0.98)
    para(t2, "Viper is open source and pre-alpha. Clone it, run the demos, "
             "and tell me what breaks.", size=20, color=MUTED, first=True,
         after=0, line=1.15)
    foot = textbox(s, 0.95, 6.7, 11.5, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    para(foot, [
        ("github.com/splanck/viper", {"color": CYAN, "bold": True}),
        ("    ·    ", {"color": BORDER}),
        ("GPL-3.0", {"color": MUTED}),
        ("    ·    ", {"color": BORDER}),
        ("Linux · macOS · Windows", {"color": MUTED}),
    ], size=16, first=True, after=0)


def add_page_numbers():
    """Tiny muted 'NN / total' top-right on every content slide."""
    total = len(prs.slides._sldIdLst)
    for idx, slide in enumerate(prs.slides, start=1):
        if idx in NO_FOOTER:
            continue
        tf = textbox(slide, SW - 1.7, 0.46, 1.0, 0.3)
        para(tf, "%02d / %d" % (idx, total), size=10.5, color=MUTED,
             align=PP_ALIGN.RIGHT, first=True, after=0)


# --------------------------------------------------------------------------- #
def main():
    slide_title()
    slide_origin()
    slide_what_is()
    slide_why()
    slide_core_idea()
    slide_languages()
    slide_same_program()
    slide_il()
    slide_architecture()
    slide_execution_vm()
    slide_execution_native()
    slide_runtime()
    slide_by_numbers()
    slide_enough()
    slide_showcase()
    slide_tooling()
    slide_thanks()
    add_page_numbers()
    prs.save(OUT)
    print("wrote", OUT, "(%d slides)" % len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
